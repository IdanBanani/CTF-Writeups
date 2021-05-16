from collections import namedtuple
import pptx
import re

# TARGET_TEXT_REGEX = re.compile("(.),\s+([a-zA-Z0-9]+\.pptx),\s+(\d+)")
TARGET_TEXT_REGEX = re.compile("^(.),\s+([a-zA-Z0-9]+\.pptx),\s+(\d+)$")

Result = namedtuple('Result', 'letter file_name slide_num')

############################################################
def get_text(ppt, slide_num):

    shapes = ppt.slides[slide_num - 1].shapes
    for shape in shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                slide_txt = run.text
                match = TARGET_TEXT_REGEX.match(run.text)
                if match is not None:
                    return Result(match.group(1), match.group(2), int(match.group(3)))
    return None

########################################################################################

flag = ""

result = Result("", "START.pptx", 1)

while True:
    try:
        result = get_text(pptx.Presentation("real_key/{}".format(result.file_name)), result.slide_num)
        print (result)
        flag += result.letter

    except pptx.exc.PackageNotFoundError:
        break

print ("Flag: {}".format(flag))