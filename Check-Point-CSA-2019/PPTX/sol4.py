import pptx
from pathlib import Path

base_dir = '.'
keys = []

def extract_key(presentation, slide_num):
    cur_key = ''
    while Path(base_dir + presentation).is_file():
        current_presentation = pptx.Presentation(base_dir + presentation)
        slide_id = current_presentation.slides._element.sldId_lst[slide_num].id
        cur_slide = current_presentation.slides.get(slide_id)

        # parse slide (***online imported code***)
        slide_text = ''
        for shape in cur_slide.shapes:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    slide_text += run.text
                break
            break

        key_char, next_presentation, slide_num = slide_text.split(',')
        presentation = next_presentation.replace(' ', '')
        slide_num = int(slide_num.replace(' ', '')) - 1
        cur_key += key_char
    return cur_key

start_presentation = pptx.Presentation(base_dir + 'START.pptx')
# parse_slide(start_presentation)
slides_counter = 0
for slide in start_presentation.slides:
    key = extract_key('START.pptx', slides_counter)
    keys.append(key)
    slides_counter += 1

counter = 1
for key in keys:
    print "Key {0}: {1}".format(counter, key)
    counter += 1

