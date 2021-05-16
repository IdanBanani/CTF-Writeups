from pptx import Presentation
import os

flag = ''

p = Presentation('START.pptx')
slide = p.slides[0]

while True:
    content = slide.shapes.title.text
    splitted = content.split(', ')

    char = splitted[0]
    next_pptx = splitted[1]
    next_pptx_slide = int(splitted[2])

    flag = flag + char

    if not os.path.exists(next_pptx):
        break

    p = Presentation(next_pptx)        
    slide = p.slides[next_pptx_slide - 1]

print(flag)