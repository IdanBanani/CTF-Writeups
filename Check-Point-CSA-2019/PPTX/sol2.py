from pptx import Presentation
import os

BASE_PATH = 'real_key/'
BASE_PATH = 'real_key/'

flag = ''
p = Presentation(BASE_PATH + 'START.pptx')
slide = p.slides[0]

while True:
    content = slide.shapes.title.text #TODO: check it out
    splitted = content.split(', ')

    char = splitted[0]
    next_pptx = BASE_PATH + splitted[1]
    next_pptx_slide = int(splitted[2])

    flag = flag + char

    if not os.path.exists(next_pptx):
        break

    p = Presentation(next_pptx)
    slide = p.slides[next_pptx_slide - 1]

print(flag)