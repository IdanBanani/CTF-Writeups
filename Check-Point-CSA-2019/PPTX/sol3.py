"""
Drop me inside the either real_key or small_key to generate the key.
"""

import pptx

START_PRS_NUM = "1"
key = ""

start_prs = pptx.Presentation("START.pptx")
next_prs = start_prs

key_char, next_prs_string, next_prs_num = "", "", START_PRS_NUM

while next_prs:
    print(next_prs.slides[int(next_prs_num) - 1].shapes.title.text)
    key_char, next_prs_string, next_prs_num = next_prs.slides[int(next_prs_num) - 1].shapes.title.text.replace(" ",
                                                                                                               "").split(
        ",")

    key += key_char
    next_prs = pptx.Presentation(next_prs_string)

    print(key)



